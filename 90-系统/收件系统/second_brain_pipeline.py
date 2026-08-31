#!/usr/bin/env python3
"""Second-brain intake, searchable extraction, queue rendering, and health checks."""

import argparse
import csv
import hashlib
import json
import os
import re
import shutil
import subprocess
import tempfile
from datetime import datetime, timedelta
from pathlib import Path


VAULT = Path("/Users/mac/Documents/Obsidian Vault")
STATE = VAULT / "90-系统/收件系统/收件状态.json"
QUEUE = VAULT / "90-系统/收件系统/收件队列.json"
QUEUE_NOTE = VAULT / "00-知识库入口/待整理资料.md"
RUN_STATE = VAULT / "90-系统/收件系统/运行状态.json"
RUN_REPORT = VAULT / "90-系统/收件系统/最近运行报告.md"
MOBILE_STATE = VAULT / "90-系统/收件系统/移动端投递状态.json"
ATTACHMENTS = VAULT / "40-资料源/附件/待整理"
SOURCE_NOTES = VAULT / "40-资料源/待整理资料页"
OCR_SCRIPT = Path("/Users/mac/.codex/skills/obsidian-vault-manager/scripts/vision_ocr.swift")
TRANSCRIBE_BIN = Path("/Users/mac/.codex/runtime/second-brain/bin/whisper")
TRANSCRIBE_PATH = "/Users/mac/.codex/runtime/second-brain/bin"
WHISPER_MODEL = Path("/Users/mac/.cache/whisper/small.pt")
SOURCE_PIPELINE = VAULT / "90-系统/收件系统/second_brain_pipeline.py"
SOURCE_OCR = VAULT / "90-系统/收件系统/vision_ocr.swift"
LAUNCH_AGENT = Path("/Users/mac/Library/LaunchAgents/com.personal.second-brain.plist")
ROOTS = [Path("/Users/mac/Desktop"), Path("/Users/mac/Documents"), Path("/Users/mac/Downloads")]
MOBILE_BRANCH = "windows-inbox"
MOBILE_INBOX = "00-知识库入口/移动端投递"
MOBILE_DROP_ROOT = Path("/Users/mac/Documents/知识库投递/Windows移动端")

TEXT_EXTS = {".md", ".txt", ".csv"}
OFFICE_EXTS = {".docx", ".pdf", ".pptx", ".xlsx"}
LEGACY_OFFICE_EXTS = {".doc", ".ppt", ".xls"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".heic"}
AUDIO_EXTS = {".mp3", ".m4a", ".wav", ".aac", ".flac", ".aiff", ".aif"}
VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".webm"}
SUPPORTED = TEXT_EXTS | OFFICE_EXTS | LEGACY_OFFICE_EXTS | IMAGE_EXTS | AUDIO_EXTS | VIDEO_EXTS

EXCLUDED_DIRS = {
    ".git", ".obsidian", ".venv", "venv", ".tox", "node_modules", "__pycache__", "qa",
    "Obsidian Vault", "Obsidian Vault整理备份", "docx_render_check", "docx_pdf_check",
}
EXCLUDED_PREFIXES = ("render", "qa_", "客户软件安装包_", "U盘备份_格式化前_")
EXCLUDED_FILES = {".DS_Store", "AGENTS.md", "README.md", "知识库投递说明.md"}
FINAL_HINTS = (
    "最终", "成品", "交付", "发布版", "专业版", "优化版", "方案", "手册", "报告", "清单",
    "逐字稿", "转写", "transcript", "课程", "访谈", "会议", "录音", "语音", "海报", "封面",
)
PRIVATE_PATTERNS = re.compile(
    r"(password|passwd|token|cookie|secret|api[_ -]?key|密码|密钥|验证码|登录信息|账号(?:名单|信息)|"
    r"vpn.*(?:分配|凭证|安装文件校验值|使用说明)|管理员(?:保管|发送前必读)|可直接发送_每人一份)", re.I,
)


def parse_args():
    parser = argparse.ArgumentParser(description="Operate the personal second-brain pipeline.")
    sub = parser.add_subparsers(dest="command", required=True)
    intake = sub.add_parser("intake")
    intake.add_argument("--since")
    intake.add_argument("--until")
    intake.add_argument("--dry-run", action="store_true")
    intake.add_argument("--max-copy-mb", type=int, default=100)
    sub.add_parser("process")
    sub.add_parser("health")
    daily = sub.add_parser("run")
    daily.add_argument("--since")
    daily.add_argument("--until")
    daily.add_argument("--max-copy-mb", type=int, default=100)
    daily.set_defaults(dry_run=False)
    sessions = sub.add_parser("session-sources")
    sessions.add_argument("--date", default=datetime.now().strftime("%Y-%m-%d"))
    return parser.parse_args()


def parse_time(value):
    dt = datetime.fromisoformat(value)
    if dt.tzinfo:
        dt = dt.astimezone().replace(tzinfo=None)
    return dt


def load_json(path, default):
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, ValueError, json.JSONDecodeError):
        return default


def save_json(path, value):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def file_hash(path):
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def is_excluded_dir(path):
    if path.name in EXCLUDED_DIRS or any(path.name.startswith(p) for p in EXCLUDED_PREFIXES):
        return True
    return (path / ".git").exists()


def is_candidate(path, start, end):
    if not path.is_file() or path.name in EXCLUDED_FILES or path.name.startswith((".", "~$")):
        return False, "系统或隐藏文件"
    try:
        path.relative_to(VAULT)
        return False, "知识库自身"
    except ValueError:
        pass
    if any(is_excluded_dir(parent) for parent in path.parents):
        return False, "排除目录"
    if path.suffix.lower() not in SUPPORTED:
        return False, "不支持格式"
    if PRIVATE_PATTERNS.search(str(path)):
        return False, "疑似敏感文件"
    stat = path.stat()
    if stat.st_size == 0:
        return False, "空文件"
    modified = datetime.fromtimestamp(stat.st_mtime)
    if not start <= modified < end:
        return False, "窗口外"
    return True, "候选"


def trust_level(path):
    parts = set(path.parts)
    lower = path.name.lower()
    if "知识库投递" in parts:
        return "自动入库"
    if str(path).startswith("/Users/mac/Documents/Codex/") and "outputs" in parts:
        return "自动入库"
    if any(h.lower() in lower for h in FINAL_HINTS):
        return "自动入库"
    return "待确认"


def kind_for(ext):
    if ext in AUDIO_EXTS:
        return "音频"
    if ext in VIDEO_EXTS:
        return "视频"
    if ext in IMAGE_EXTS:
        return "图片"
    if ext in OFFICE_EXTS | LEGACY_OFFICE_EXTS:
        return "文档"
    return "文本"


def next_action(ext, copied):
    if not copied:
        return "确认是否入库"
    if ext in AUDIO_EXTS:
        return "待转写和讲者确认"
    if ext in VIDEO_EXTS:
        return "待建立稳定引用和文本说明"
    if ext in IMAGE_EXTS:
        return "待OCR或人工说明"
    if ext in LEGACY_OFFICE_EXTS:
        return "待转换为新版格式后提取文本"
    return "自动提取可搜索文本"


def unique_target(source):
    target = ATTACHMENTS / source.name
    if not target.exists():
        return target
    source_hash = file_hash(source)
    if file_hash(target) == source_hash:
        return target
    for i in range(2, 1000):
        alt = ATTACHMENTS / f"{source.stem}-{i}{source.suffix}"
        if not alt.exists() or file_hash(alt) == source_hash:
            return alt
    raise RuntimeError(f"cannot allocate target for {source}")


def entry_id(path, digest):
    return digest[:12] if digest else hashlib.sha256(str(path).encode()).hexdigest()[:12]


def ingest_mobile_inbox():
    """Stage new, supported blobs from the Windows-only branch into the trusted Mac intake folder."""
    result = {"branch": MOBILE_BRANCH, "seen": 0, "staged": 0, "skipped": 0, "failures": []}
    fetch = subprocess.run(
        ["git", "fetch", "--quiet", "origin", f"refs/heads/{MOBILE_BRANCH}:refs/remotes/origin/{MOBILE_BRANCH}"],
        cwd=VAULT, capture_output=True, text=True, timeout=120,
    )
    if fetch.returncode != 0:
        result["failures"].append((fetch.stderr or fetch.stdout or "无法拉取移动端投递分支").strip()[-1000:])
        return result
    listed = subprocess.run(
        ["git", "ls-tree", "-r", "--name-only", f"origin/{MOBILE_BRANCH}", "--", MOBILE_INBOX],
        cwd=VAULT, capture_output=True, text=True, timeout=60,
    )
    if listed.returncode != 0:
        result["failures"].append((listed.stderr or listed.stdout or "无法读取移动端投递目录").strip()[-1000:])
        return result
    state = load_json(MOBILE_STATE, {})
    delivered = set(state.get("delivered_sha256", []))
    for rel in (line.strip() for line in listed.stdout.splitlines()):
        if not rel:
            continue
        name = Path(rel).name
        extension = Path(name).suffix.lower()
        result["seen"] += 1
        if name in EXCLUDED_FILES or name.startswith(".") or extension not in SUPPORTED or PRIVATE_PATTERNS.search(name):
            result["skipped"] += 1
            continue
        blob = subprocess.run(["git", "show", f"origin/{MOBILE_BRANCH}:{rel}"], cwd=VAULT,
                              capture_output=True, timeout=120)
        if blob.returncode != 0:
            result["failures"].append(f"无法读取 {rel}")
            continue
        digest = hashlib.sha256(blob.stdout).hexdigest()
        if not blob.stdout or digest in delivered:
            result["skipped"] += 1
            continue
        MOBILE_DROP_ROOT.mkdir(parents=True, exist_ok=True)
        target = MOBILE_DROP_ROOT / f"{digest[:12]}-{name}"
        if not target.exists():
            target.write_bytes(blob.stdout)
        delivered.add(digest)
        result["staged"] += 1
    state["delivered_sha256"] = sorted(delivered)
    state["last_checked_at"] = datetime.now().isoformat(timespec="seconds")
    save_json(MOBILE_STATE, state)
    return result


def initialize_existing(queue):
    known_paths = {item.get("vault_path") for item in queue}
    if not ATTACHMENTS.exists():
        return queue
    for path in sorted(ATTACHMENTS.iterdir()):
        if not path.is_file():
            continue
        rel = path.relative_to(VAULT).as_posix()
        if rel in known_paths:
            continue
        digest = file_hash(path)
        queue.append({
            "id": entry_id(path, digest), "received_at": datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec="seconds"),
            "name": path.name, "kind": kind_for(path.suffix.lower()), "extension": path.suffix.lower(),
            "source_path": "历史收件来源见待整理资料旧记录", "vault_path": rel, "sha256": digest,
            "trust": "历史已收件", "status": "待处理", "next_action": next_action(path.suffix.lower(), True),
            "note_path": "", "size_mb": round(path.stat().st_size / 1024 / 1024, 2), "related_project": "待判断",
        })
    return queue


def render_queue(queue):
    order = {"待处理": 0, "来源缺失": 1, "处理失败": 2, "待确认": 3, "仅登记": 4, "待转写": 5, "待OCR": 6,
             "已转写待复核": 7, "已OCR待复核": 8, "已结构化": 9, "不入库": 10}
    queue = sorted(queue, key=lambda x: (order.get(x.get("status"), 8), x.get("received_at", "")), reverse=False)
    counts = {}
    for item in queue:
        counts[item["status"]] = counts.get(item["status"], 0) + 1
    summary = "；".join(f"{k}{v}项" for k, v in sorted(counts.items())) or "当前没有待处理资料"
    lines = [
        "---", "title: 待整理资料", "type: system", "status: reviewed",
        "summary: 每日新增资料的唯一人工处理入口。", f"updated: {datetime.now().strftime('%Y-%m-%d')}",
        "topics:", "  - 收件", "  - 待整理", "related:", '  - "[[90-系统/收件系统/收件系统说明]]"', "---", "",
        "# 待整理资料", "", "> [!important]", "> 本页由收件队列生成。原件不因入库而删除；未形成可搜索文本的资料不算完成知识化。", "",
        "## 当前状态", "", f"{summary}。", "", "## 处理队列", "",
        "| 收件时间 | 资料 | 类型 | 来源级别 | 状态 | 下一步 | 关联项目 |", "| --- | --- | --- | --- | --- | --- | --- |",
    ]
    for item in queue:
        if item.get("note_path"):
            label = f"[[{item['note_path']}|{item['name']}]]"
        elif item.get("vault_path"):
            label = f"[[{item['vault_path']}|{item['name']}]]"
        else:
            label = f"`{item['name']}`"
        lines.append("| {received} | {label} | {kind} | {trust} | {status} | {action} | {project} |".format(
            received=item.get("received_at", "")[:16].replace("T", " "), label=label, kind=item.get("kind", ""),
            trust=item.get("trust", ""), status=item.get("status", ""), action=item.get("next_action", ""),
            project=item.get("related_project", "待判断"),
        ))
    lines += [
        "", "## 处理原则", "", "1. 先确认来源、隐私边界和关联项目。", "2. 音频先转写，图片先OCR或补说明，办公文件先提取文本。",
        "3. 资料页只记录来源事实；可复用结论另行进入知识层。", "4. 项目证据进入项目层，稳定步骤进入技能与SOP层。",
        "5. `待确认`项目只登记，不自动复制；确认无长期价值后标记`不入库`。", "",
        "[[00-知识库入口/知识库总览|返回知识底座]]", "",
    ]
    QUEUE_NOTE.write_text("\n".join(lines), encoding="utf-8")


def intake(args):
    state = load_json(STATE, {})
    start = parse_time(args.since) if args.since else parse_time(state.get("last_success_at", (datetime.now() - timedelta(days=1)).isoformat()))
    end = parse_time(args.until) if args.until else datetime.now()
    if start >= end:
        raise SystemExit("scan start must be before end")
    queue = initialize_existing(load_json(QUEUE, []))
    original_queue_size = len(queue)
    queue = [item for item in queue if not PRIVATE_PATTERNS.search(
        f"{item.get('name', '')} {item.get('source_path', '')}"
    )]
    sensitive_queue_removed = original_queue_size - len(queue)
    by_hash = {item.get("sha256"): item for item in queue if item.get("sha256")}
    by_source = {item.get("source_path"): item for item in queue}
    candidates = []
    skipped = {}
    for root in ROOTS:
        if not root.exists():
            continue
        for current, dirs, files in os.walk(root):
            current_path = Path(current)
            dirs[:] = [d for d in dirs if not is_excluded_dir(current_path / d)]
            for name in files:
                path = current_path / name
                ok, reason = is_candidate(path, start, end)
                if ok:
                    candidates.append(path)
                else:
                    skipped[reason] = skipped.get(reason, 0) + 1
    added = []
    for source in sorted(set(candidates), key=lambda p: str(p)):
        if str(source) in by_source:
            continue
        digest = file_hash(source)
        if digest in by_hash:
            continue
        ext = source.suffix.lower()
        trust = trust_level(source)
        size_mb = source.stat().st_size / 1024 / 1024
        copied = trust == "自动入库" and ext not in VIDEO_EXTS and size_mb <= args.max_copy_mb
        target = unique_target(source) if copied else None
        if copied and not args.dry_run:
            ATTACHMENTS.mkdir(parents=True, exist_ok=True)
            if not target.exists():
                shutil.copy2(source, target)
        item = {
            "id": entry_id(source, digest), "received_at": datetime.now().isoformat(timespec="seconds"), "name": source.name,
            "kind": kind_for(ext), "extension": ext, "source_path": str(source),
            "vault_path": target.relative_to(VAULT).as_posix() if target else "", "sha256": digest, "trust": trust,
            "status": "待处理" if copied else ("仅登记" if trust == "自动入库" else "待确认"),
            "next_action": next_action(ext, copied), "note_path": "", "size_mb": round(size_mb, 2), "related_project": "待判断",
        }
        added.append(item)
        by_hash[digest] = item
    if not args.dry_run:
        queue.extend(added)
        save_json(QUEUE, queue)
        save_json(STATE, {"last_success_at": datetime.now().isoformat(timespec="seconds"), "updated_by": "second_brain_pipeline.py"})
        render_queue(queue)
    print(f"scan_start={start.isoformat(timespec='seconds')}")
    print(f"scan_end={end.isoformat(timespec='seconds')}")
    print(f"candidates={len(candidates)}")
    print(f"new_queue_items={len(added)}")
    print(f"sensitive_queue_removed={sensitive_queue_removed}")
    for item in added:
        print(f"{item['status']}\t{item['kind']}\t{item['source_path']}")
    print("skipped=" + json.dumps(skipped, ensure_ascii=False, sort_keys=True))
    return {"start": start.isoformat(timespec="seconds"), "end": end.isoformat(timespec="seconds"),
            "candidates": len(candidates), "added": len(added), "skipped": skipped,
            "items": [report_item(item) for item in added],
            "sensitive_queue_removed": sensitive_queue_removed}


def clean_text(value):
    return value.replace("\x00", "").strip()


def extract_text(path):
    ext = path.suffix.lower()
    if ext in {".md", ".txt"}:
        return path.read_text(encoding="utf-8", errors="replace")
    if ext == ".csv":
        return path.read_text(encoding="utf-8", errors="replace")
    if ext == ".docx":
        from docx import Document
        doc = Document(path)
        chunks = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n\n".join(chunks)
    if ext == ".pdf":
        from pypdf import PdfReader
        return "\n\n".join((page.extract_text() or "") for page in PdfReader(path).pages)
    if ext == ".pptx":
        from pptx import Presentation
        chunks = []
        for number, slide in enumerate(Presentation(path).slides, 1):
            text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            chunks.append(f"## 第{number}页\n\n" + "\n\n".join(text))
        return "\n\n".join(chunks)
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        chunks = []
        for sheet in wb.worksheets:
            chunks.append(f"## 工作表：{sheet.title}")
            for index, row in enumerate(sheet.iter_rows(values_only=True), 1):
                if index > 1000:
                    chunks.append("[仅提取前1000行]")
                    break
                chunks.append(" | ".join("" if value is None else str(value) for value in row))
        return "\n".join(chunks)
    raise ValueError("format requires conversion or transcription")


def ocr_image(path):
    if not OCR_SCRIPT.exists():
        raise RuntimeError("macOS Vision OCR脚本不存在")
    result = subprocess.run(["swift", str(OCR_SCRIPT), str(path)], capture_output=True, text=True, timeout=1800)
    if result.returncode != 0:
        raise RuntimeError((result.stderr or "Vision OCR失败").strip())
    text = clean_text(result.stdout)
    if not text:
        raise RuntimeError("Vision OCR未识别到文字")
    return text


def format_timestamp(seconds):
    seconds = max(0, int(float(seconds)))
    hours, remain = divmod(seconds, 3600)
    minutes, seconds = divmod(remain, 60)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}"


def transcribe_audio(path):
    if not TRANSCRIBE_BIN.exists():
        raise RuntimeError("本地Whisper运行环境不存在")
    env = os.environ.copy()
    env["PATH"] = TRANSCRIBE_PATH + os.pathsep + env.get("PATH", "")
    with tempfile.TemporaryDirectory(prefix="second-brain-whisper-") as temp_dir:
        command = [str(TRANSCRIBE_BIN), str(path), "--model", "small", "--device", "cpu",
                   "--output_dir", temp_dir, "--output_format", "json", "--verbose", "False"]
        result = subprocess.run(command, capture_output=True, text=True, env=env, timeout=14400)
        outputs = sorted(Path(temp_dir).glob("*.json"))
        if result.returncode != 0 or not outputs:
            detail = (result.stderr or result.stdout or "Whisper转写失败").strip()
            raise RuntimeError(detail[-1000:])
        payload = load_json(outputs[0], {})
        full_text = clean_text(payload.get("text", ""))
        if not full_text:
            raise RuntimeError("Whisper未识别到语音文字")
        timeline = []
        for segment in payload.get("segments", []):
            text = clean_text(segment.get("text", ""))
            if text:
                timeline.append(f"- [{format_timestamp(segment.get('start', 0))}–{format_timestamp(segment.get('end', 0))}] {text}")
        parts = ["### 完整转写", "", full_text]
        if timeline:
            parts += ["", "### 分段时间轴", "", *timeline]
        return "\n".join(parts)


def safe_title(name):
    title = re.sub(r"[\\/:*?\"<>|#\[\]]", "-", Path(name).stem).strip(" .-")
    return title[:80] or "未命名资料"


def create_source_note(item, text=""):
    SOURCE_NOTES.mkdir(parents=True, exist_ok=True)
    title = safe_title(item["name"])
    note = SOURCE_NOTES / f"{item['id']}-{title}.md"
    source_link = f"[[{item['vault_path']}|原始文件]]" if item.get("vault_path") else f"`{item['source_path']}`"
    searchable = bool(text.strip())
    processing = item.get("processing_method") or ("已提取文本" if searchable else item["next_action"])
    body = [
        "---", f'title: "{title}"', "type: source", "status: raw", f'summary: "自动收件的{item["kind"]}资料，等待人工确认来源边界和知识价值。"',
        f"updated: {datetime.now().strftime('%Y-%m-%d')}", "topics:", "  - 待整理", f"source_file: \"{item.get('vault_path') or item['source_path']}\"",
        f"searchable: {'true' if searchable else 'false'}", f'processing: "{processing}"', "related: []", "---", "", f"# {title}", "",
        "> [!warning] 来源边界", "> 本页由收件系统自动生成，只表示资料已经保存或登记，不代表其中观点已经核验，也不代表用户立场。", "",
        "## 资料信息", "", f"- 类型：{item['kind']}", f"- 原始文件：{source_link}", f"- 原始来源：`{item['source_path']}`",
        f"- 收件时间：{item['received_at']}", f"- 当前处理：{processing}", "", "## 核心主线", "", "待人工或AI通读后补充。", "",
    ]
    if searchable:
        body += ["## 自动提取文本", "", clean_text(text), ""]
    else:
        body += ["## 待处理", "", f"- [ ] {item['next_action']}"]
        if item.get("last_error"):
            body += [f"- 最近错误：`{item['last_error']}`"]
        body += ["- [ ] 确认讲者、来源、时间和隐私边界", "- [ ] 形成可搜索文本后再提炼知识", ""]
    body += ["## 六要素提炼", "", "### 认知 / 观点", "", "### 方法论", "", "### 故事 / 数据", "", "### 情绪表达", "", "### 可发展选题", "", "### 待办", "", "## 分流结果", "", "- 知识：待判断", "- 项目证据：待判断", "- SOP或Skill：待判断", ""]
    note.write_text("\n".join(body), encoding="utf-8")
    return note.relative_to(VAULT).with_suffix("").as_posix(), searchable


def report_item(item):
    return {key: item.get(key, "") for key in (
        "name", "kind", "source_path", "vault_path", "note_path", "trust", "status",
        "next_action", "processing_method", "related_project",
    )}


def process_queue():
    queue = initialize_existing(load_json(QUEUE, []))
    changed = 0
    failures = []
    processed_items = []
    processed_by_type = {"文档文本": 0, "图片OCR": 0, "音频转写": 0, "待转换": 0}
    for item in queue:
        if item.get("note_path") or not item.get("vault_path"):
            continue
        path = VAULT / item["vault_path"]
        if not path.exists():
            item["status"] = "来源缺失"
            item["next_action"] = "核验原始文件和Vault副本"
            changed += 1
            processed_items.append(report_item(item))
            continue
        try:
            if item["extension"] in AUDIO_EXTS:
                text = transcribe_audio(path)
                item["processing_method"] = "本地Whisper small自动转写，待人工复核"
            elif item["extension"] in IMAGE_EXTS:
                text = ocr_image(path)
                item["processing_method"] = "macOS Vision自动OCR，待人工复核"
            else:
                text = extract_text(path)
                item["processing_method"] = "已自动提取文本"
            item.pop("last_error", None)
        except Exception as error:
            text = ""
            item["last_error"] = str(error).replace("\n", " ")[:500]
            failures.append({"name": item["name"], "error": item["last_error"]})
        note_path, searchable = create_source_note(item, text)
        item["note_path"] = note_path
        if searchable and item["extension"] in AUDIO_EXTS:
            item["status"] = "已转写待复核"
            item["next_action"] = "人工复核讲者、专有名词和关键数字"
            processed_by_type["音频转写"] += 1
        elif searchable and item["extension"] in IMAGE_EXTS:
            item["status"] = "已OCR待复核"
            item["next_action"] = "人工对照图片复核错字和版面顺序"
            processed_by_type["图片OCR"] += 1
        elif searchable:
            item["status"] = "已结构化"
            item["next_action"] = "人工通读并决定是否提炼知识"
            processed_by_type["文档文本"] += 1
        elif item["extension"] in AUDIO_EXTS:
            item["status"] = "待转写"
        elif item["extension"] in IMAGE_EXTS:
            item["status"] = "待OCR"
        else:
            item["status"] = "待转换"
            processed_by_type["待转换"] += 1
        changed += 1
        processed_items.append(report_item(item))
    save_json(QUEUE, queue)
    render_queue(queue)
    print(f"processed={changed}")
    print(f"processing_failures={len(failures)}")
    print("processed_by_type=" + json.dumps(processed_by_type, ensure_ascii=False, sort_keys=True))
    for failure in failures:
        print(f"FAILED\t{failure['name']}\t{failure['error']}")
    return {"processed": changed, "processed_by_type": processed_by_type, "failures": failures,
            "items": processed_items}


def health():
    queue = load_json(QUEUE, [])
    counts = {}
    for item in queue:
        counts[item.get("status", "未知")] = counts.get(item.get("status", "未知"), 0) + 1
    last_review = None
    review_root = VAULT / "90-系统/04-Codex永久记忆/00-每日复盘"
    reviews = sorted(p for p in review_root.rglob("*.md") if p.name != "索引.md")
    if reviews:
        last_review = reviews[-1]
    files = [p for p in VAULT.rglob("*") if p.is_file()]
    rels = {p.relative_to(VAULT).as_posix() for p in files}
    stems = {p.relative_to(VAULT).with_suffix("").as_posix() for p in files}
    broken = []
    link_re = re.compile(r"\[\[([^\]]+)\]\]")
    for note in VAULT.rglob("*.md"):
        if note.relative_to(VAULT).parts[0] == "99-归档":
            continue
        for match in link_re.finditer(note.read_text(encoding="utf-8", errors="ignore")):
            raw = match.group(1).split("|", 1)[0].split("#", 1)[0].strip().lstrip("/")
            if not raw:
                continue
            candidates = [raw, raw + ".md", raw + ".base", (note.parent / raw).relative_to(VAULT).as_posix()]
            if not any(c in rels or c in stems for c in candidates):
                broken.append((note.relative_to(VAULT).as_posix(), raw))
    print("queue=" + json.dumps(counts, ensure_ascii=False, sort_keys=True))
    print(f"last_intake={load_json(STATE, {}).get('last_success_at', 'unknown')}")
    print(f"last_review={last_review.relative_to(VAULT) if last_review else 'none'}")
    print(f"active_broken_links={len(broken)}")
    for note, target in broken[:30]:
        print(f"BROKEN\t{note}\t{target}")
    now = datetime.now()
    expected_review = (now - timedelta(days=1)).date() if now.hour < 23 else now.date()
    latest_review_date = None
    if last_review:
        match = re.match(r"(\d{4}-\d{2}-\d{2})", last_review.name)
        if match:
            latest_review_date = datetime.strptime(match.group(1), "%Y-%m-%d").date()
    review_lag = (expected_review - latest_review_date).days if latest_review_date else 9999
    pending_statuses = {"待处理", "来源缺失", "处理失败", "待确认", "待转写", "待OCR", "待转换"}
    old_pending = []
    for item in queue:
        if item.get("status") not in pending_statuses:
            continue
        try:
            received = parse_time(item["received_at"])
        except Exception:
            continue
        if now - received > timedelta(days=7):
            old_pending.append(item["name"])
    launch_result = subprocess.run(["launchctl", "print", f"gui/{os.getuid()}/com.personal.second-brain"],
                                   capture_output=True, text=True, timeout=30)
    launch_loaded = launch_result.returncode == 0
    launch_exit_match = re.search(r"last exit code = ([^\n]+)", launch_result.stdout) if launch_loaded else None
    launch_last_exit = launch_exit_match.group(1).strip() if launch_exit_match else "unknown"
    runtime_checks = {
        "pipeline_runtime": Path(__file__).exists(), "ocr_runtime": OCR_SCRIPT.exists(),
        "transcribe_runtime": TRANSCRIBE_BIN.exists(), "whisper_model": WHISPER_MODEL.exists(),
        "launch_agent_file": LAUNCH_AGENT.exists(), "launch_agent_loaded": launch_loaded,
        "pipeline_copy_match": SOURCE_PIPELINE.exists() and file_hash(SOURCE_PIPELINE) == file_hash(Path(__file__)),
        "ocr_copy_match": SOURCE_OCR.exists() and OCR_SCRIPT.exists() and file_hash(SOURCE_OCR) == file_hash(OCR_SCRIPT),
    }
    runtime_failures = [name for name, ok in runtime_checks.items() if not ok]
    git_result = subprocess.run(["git", "status", "--porcelain"], cwd=VAULT, capture_output=True, text=True, timeout=30)
    git_changes = len([line for line in git_result.stdout.splitlines() if line.strip()]) if git_result.returncode == 0 else -1
    print(f"review_lag_days={max(0, review_lag)}")
    print(f"old_pending_over_7d={len(old_pending)}")
    print(f"runtime_failures={len(runtime_failures)}")
    print(f"launch_agent_last_exit={launch_last_exit}")
    print(f"git_worktree_changes={git_changes}")
    for name in runtime_failures:
        print(f"RUNTIME_FAILED\t{name}")
    return {"queue": counts, "last_intake": load_json(STATE, {}).get("last_success_at", "unknown"),
            "last_review": str(last_review.relative_to(VAULT)) if last_review else "none", "broken": broken,
            "review_lag": max(0, review_lag), "old_pending": old_pending, "runtime_failures": runtime_failures,
            "launch_agent_last_exit": launch_last_exit, "git_changes": git_changes}


def write_run_report(started_at, mobile_result, intake_result, process_result, health_result):
    finished_at = datetime.now().isoformat(timespec="seconds")
    hard_failure = bool(mobile_result["failures"] or process_result["failures"] or
                        health_result["broken"] or health_result["runtime_failures"])
    needs_attention = bool(health_result["review_lag"] or health_result["old_pending"])
    status = "失败" if hard_failure else ("需要处理" if needs_attention else "成功")
    payload = {"status": status, "started_at": started_at, "finished_at": finished_at,
               "mobile_inbox": mobile_result, "intake": intake_result, "process": process_result, "health": {
                   "queue": health_result["queue"], "last_intake": health_result["last_intake"],
                   "last_review": health_result["last_review"], "review_lag_days": health_result["review_lag"],
                   "broken_count": len(health_result["broken"]), "old_pending": health_result["old_pending"],
                   "launch_agent_last_exit": health_result["launch_agent_last_exit"],
                   "runtime_failures": health_result["runtime_failures"], "git_worktree_changes": health_result["git_changes"],
               }}
    save_json(RUN_STATE, payload)
    queue_summary = "；".join(f"{key}{value}项" for key, value in sorted(health_result["queue"].items())) or "无"
    lines = ["---", "title: 最近运行报告", "type: system", "status: reviewed",
             f"updated: {datetime.now().strftime('%Y-%m-%d')}", "---", "", "# 最近运行报告", "",
             f"- 运行状态：**{status}**", f"- 开始：{started_at}", f"- 完成：{finished_at}",
             f"- Windows移动端投递：发现{mobile_result['seen']}项，暂存{mobile_result['staged']}项，跳过{mobile_result['skipped']}项",
             f"- 扫描窗口：{intake_result['start']} → {intake_result['end']}",
             f"- 新增队列：{intake_result['added']}项", f"- 本次处理：{process_result['processed']}项",
             f"- 凭证类资料安全隔离：{intake_result.get('sensitive_queue_removed', 0)}项（只移除队列登记，不删除源文件）",
             f"- 文档文本：{process_result['processed_by_type']['文档文本']}项",
             f"- 图片OCR：{process_result['processed_by_type']['图片OCR']}项",
             f"- 音频转写：{process_result['processed_by_type']['音频转写']}项",
             f"- 队列：{queue_summary}", f"- 活动区断链：{len(health_result['broken'])}处",
             f"- LaunchAgent运行前上次退出码：{health_result['launch_agent_last_exit']}",
             f"- Git同步前变更：{health_result['git_changes']}项", ""]
    if intake_result.get("items"):
        lines += ["## 本次发现的具体资料", ""]
        for item in intake_result["items"]:
            destination = f"[[{item['vault_path']}|Vault原件]]" if item.get("vault_path") else "仅登记来源，未复制入库"
            lines += [
                f"### {item['name']}", "",
                f"- 类型：{item['kind']}",
                f"- 来源：`{item['source_path']}`",
                f"- 来源级别：{item['trust']}",
                f"- 当前状态：{item['status']}",
                f"- 保存结果：{destination}",
                f"- 下一步：{item['next_action']}", "",
            ]
    else:
        lines += ["## 本次发现的具体资料", "", "- 本扫描窗口没有发现新的合格资料。", ""]
    if process_result.get("items"):
        lines += ["## 本次实际完成的整理", ""]
        for item in process_result["items"]:
            note = f"[[{item['note_path']}|查看整理后的资料页]]" if item.get("note_path") else "未生成资料页"
            lines += [
                f"### {item['name']}", "",
                f"- 整理方式：{item.get('processing_method') or item['next_action']}",
                f"- 整理后状态：{item['status']}",
                f"- 整理结果：{note}",
                f"- 下一步：{item['next_action']}", "",
            ]
    else:
        lines += ["## 本次实际完成的整理", "", "- 本次没有完成新的全文提取、OCR或音频转写。", ""]
    if mobile_result["failures"]:
        lines += ["## 移动端投递异常", ""] + [f"- {item}" for item in mobile_result["failures"]] + [""]
    if process_result["failures"]:
        lines += ["## 处理失败", ""] + [f"- `{item['name']}`：{item['error']}" for item in process_result["failures"]] + [""]
    if health_result["runtime_failures"]:
        lines += ["## 运行环境异常", ""] + [f"- `{name}`" for name in health_result["runtime_failures"]] + [""]
    if health_result["review_lag"]:
        lines += ["## 复盘积压", "", f"- 每日复盘落后{health_result['review_lag']}天。", ""]
    if health_result["old_pending"]:
        lines += ["## 超过7天的待处理资料", ""] + [f"- `{name}`" for name in health_result["old_pending"]] + [""]
    lines += ["[[00-知识库入口/待整理资料|查看待整理资料]]", ""]
    RUN_REPORT.write_text("\n".join(lines), encoding="utf-8")
    return status, hard_failure


def daily_run(args):
    started_at = datetime.now().isoformat(timespec="seconds")
    mobile_result = ingest_mobile_inbox()
    intake_result = intake(args)
    process_result = process_queue()
    health_result = health()
    status, hard_failure = write_run_report(started_at, mobile_result, intake_result, process_result, health_result)
    print(f"run_status={status}")
    if hard_failure:
        raise SystemExit(1)


def session_sources(date_text):
    roots = [Path("/Users/mac/.codex/sessions"), Path("/Users/mac/.codex/archived_sessions")]
    matches = []
    for root in roots:
        if not root.exists():
            continue
        matches.extend(p for p in root.rglob("*.jsonl") if date_text in p.name or date_text.replace("-", "/") in str(p))
    print(f"date={date_text}")
    print(f"session_files={len(matches)}")
    for path in sorted(matches):
        print(path)


def main():
    args = parse_args()
    if args.command == "intake":
        intake(args)
    elif args.command == "process":
        process_queue()
    elif args.command == "health":
        health()
    elif args.command == "run":
        daily_run(args)
    elif args.command == "session-sources":
        session_sources(args.date)


if __name__ == "__main__":
    main()
